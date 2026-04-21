#!/usr/bin/env python3
r"""Generate a daily agenda slide (.pptx) from a lesson plan's markdown sections.

Usage:
  python generate_agenda_slide.py \
    --plan path/to/plan.md \
    --date 2026-04-22 \
    --subject chem \
    --output ~/Documents/Lesson\ Plan\ Magic/outputs/2026-04-22_chem_agenda.pptx \
    [--template path/to/agenda-slide-template.pptx]

If --template is provided and the file exists, uses it as a base (layout 0,
placeholder-fill mode).  Otherwise creates a clean deck from scratch with
back-of-room-readable sizing (36pt title, 24pt body, neutral palette).

Requires python-pptx (see requirements.txt).
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

_HERE = Path(__file__).resolve().parent
if str(_HERE) not in sys.path:
    sys.path.insert(0, str(_HERE))
from artifact_common import load_config, resolve_output_path
# Hard dependency: pii_scan carries the PII checks and the plan
# loader. A broken/incomplete install must abort loudly, not fall back to
# no-op stubs that would persist student names to disk.
from pii_scan import (
    scan_for_pii,
    scan_pptx_for_pii,
    resolve_allowlist,
    load_plan_md,
    replace_pptx_placeholders,
    find_unresolved_pptx_placeholders,
)
from content_schema import (
    ContentValidationError,
    load_and_validate_content,
    validate_agenda_content,
)
from plan_parser import parse_plan


# ---------------------------------------------------------------------------
# Markdown → day fields
# ---------------------------------------------------------------------------

def extract_day_fields(md: str, target_date: str) -> dict | None:
    """Pull the target date's sections from a weekly markdown plan.

    Returns dict with keys: date, day_name, learning_intention,
    success_criteria (list), agenda (list), do_now, materials (list).
    Returns None if the target date isn't found in the plan.
    """
    week = parse_plan(md)
    for day in week.days:
        if day.date != target_date:
            continue
        return {
            "date": day.date,
            "day_name": day.day_name,
            "learning_intention": day.learning_intention,
            "success_criteria": list(day.success_criteria),
            "agenda": list(day.agenda),
            "do_now": day.do_now,
            "materials": list(day.materials),
            "subject": week.subject,
        }
    return None


# ---------------------------------------------------------------------------
# Slide generation — from scratch
# ---------------------------------------------------------------------------

def _add_textbox(slide, left, top, width, height, text, font_size,
                 bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Add a textbox with a single styled run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size,
                     numbered=False, color=None):
    """Add a textbox with bullet or numbered items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        prefix = f"{i + 1}. " if numbered else "• "
        run = p.add_run()
        run.text = f"{prefix}{item}"
        run.font.size = Pt(font_size)
        if color:
            run.font.color.rgb = color
    return txBox


def _build_from_scratch(fields: dict):
    """Build a clean agenda slide with back-of-room-readable sizing.

    Returns the in-memory Presentation so callers can scan it before saving.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    dark = RGBColor(0x2D, 0x2D, 0x2D)
    accent = RGBColor(0x1A, 0x56, 0x76)
    muted = RGBColor(0x55, 0x55, 0x55)

    margin = Inches(0.6)
    content_w = prs.slide_width - 2 * margin

    # Title: date + day name
    title_text = fields["date"]
    if fields.get("day_name"):
        title_text += f" — {fields['day_name']}"
    _add_textbox(slide, margin, Inches(0.3), content_w, Inches(0.7),
                 title_text, 36, bold=True, color=accent)

    # Subject footer
    if fields.get("subject"):
        _add_textbox(slide, margin, Inches(6.8), content_w, Inches(0.5),
                     fields["subject"], 16, color=muted, alignment=PP_ALIGN.RIGHT)

    y = Inches(1.1)

    # Learning Intention
    if fields["learning_intention"]:
        _add_textbox(slide, margin, y, content_w, Inches(0.35),
                     "Learning Intention", 20, bold=True, color=accent)
        y += Inches(0.35)
        _add_textbox(slide, margin, y, content_w, Inches(0.6),
                     fields["learning_intention"], 24, color=dark)
        y += Inches(0.65)

    # Success Criteria
    if fields["success_criteria"]:
        _add_textbox(slide, margin, y, content_w, Inches(0.35),
                     "Success Criteria", 20, bold=True, color=accent)
        y += Inches(0.35)
        _add_bullet_list(slide, margin, y, content_w, Inches(1.2),
                         fields["success_criteria"], 22, color=dark)
        y += Inches(0.35 * len(fields["success_criteria"]) + 0.3)

    # Agenda
    if fields["agenda"]:
        _add_textbox(slide, margin, y, content_w, Inches(0.35),
                     "Agenda", 20, bold=True, color=accent)
        y += Inches(0.35)
        _add_bullet_list(slide, margin, y, content_w, Inches(2.0),
                         fields["agenda"], 22, numbered=True, color=dark)
        y += Inches(0.35 * len(fields["agenda"]) + 0.3)

    # Do-Now (lower third)
    if fields.get("do_now"):
        _add_textbox(slide, margin, y, content_w, Inches(0.35),
                     "Do Now", 20, bold=True, color=accent)
        y += Inches(0.35)
        _add_textbox(slide, margin, y, content_w, Inches(0.6),
                     fields["do_now"], 22, color=dark)
        y += Inches(0.65)

    # Homework (rendered explicitly — content["homework"] populates this via --content-json).
    if fields.get("homework"):
        _add_textbox(slide, margin, y, content_w, Inches(0.35),
                     "Homework", 20, bold=True, color=accent)
        y += Inches(0.35)
        _add_textbox(slide, margin, y, content_w, Inches(0.6),
                     fields["homework"], 22, color=dark)

    return prs


# ---------------------------------------------------------------------------
# Slide generation — from template
# ---------------------------------------------------------------------------

def _build_from_template(fields: dict, template_path: Path):
    """Fill an existing .pptx template's placeholders and return the Presentation.

    Looks for text containing {{DATE}}, {{LEARNING_INTENTION}}, etc.
    in slide shapes and replaces them. The caller saves the result after
    scanning for PII.
    """
    prs = Presentation(str(template_path))

    replacements = {
        "{{DATE}}": fields["date"],
        "{{DAY_NAME}}": fields.get("day_name", ""),
        "{{SUBJECT}}": fields.get("subject", ""),
        "{{LEARNING_INTENTION}}": fields["learning_intention"],
        "{{SUCCESS_CRITERIA}}": "\n".join(f"• {c}" for c in fields["success_criteria"]),
        "{{AGENDA}}": "\n".join(f"{i+1}. {a}" for i, a in enumerate(fields["agenda"])),
        "{{DO_NOW}}": fields.get("do_now", ""),
        "{{MATERIALS}}": "\n".join(f"• {m}" for m in fields.get("materials", [])),
        "{{HOMEWORK}}": fields.get("homework", ""),
    }

    replace_pptx_placeholders(prs, replacements)

    return prs


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> int:
    ap = argparse.ArgumentParser(description="Generate agenda slide from lesson plan.")
    ap.add_argument("--plan", required=True, help="Path to markdown or .docx lesson plan")
    ap.add_argument("--date", required=True, help="Target date (YYYY-MM-DD)")
    ap.add_argument("--subject", help="Subject id (for filename / footer)")
    ap.add_argument("--output", required=True, help="Output .pptx path")
    ap.add_argument("--template", help="Optional .pptx template to fill")
    ap.add_argument("--content-file", help="Path to content JSON (agenda/homework override)")
    ap.add_argument("--content-json", help="Inline content JSON (agenda/homework override)")
    ap.add_argument("--config", help="Config file for PII allowlist (defaults to ~/Documents/Lesson Plan Magic/config.yaml)")
    ap.add_argument(
        "--allow-names",
        help=(
            "Comma-separated bare names to allow "
            "(teacher/public figures only; keyword-shaped PII still fails)"
        ),
    )
    ap.add_argument("--allow-anywhere", action="store_true", help="Allow writing outside ~/Documents/Lesson Plan Magic/outputs/")
    args = ap.parse_args()

    plan_path = Path(args.plan)
    try:
        md = load_plan_md(plan_path)
    except (FileNotFoundError, ValueError) as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    try:
        fields = extract_day_fields(md, args.date)
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    if fields is None:
        print(f"Error: date {args.date} not found in plan.", file=sys.stderr)
        return 1

    if args.subject:
        fields["subject"] = args.subject

    # Parse content overrides (if provided)
    try:
        content = load_and_validate_content(
            content_file=args.content_file,
            content_json=args.content_json,
            validator=validate_agenda_content,
        )
    except ContentValidationError as e:
        print(f"Error: invalid content JSON: {e}", file=sys.stderr)
        return 1

    # Override fields with content if provided
    if "agenda" in content:
        fields["agenda"] = content["agenda"]
    if "homework" in content:
        fields["homework"] = content["homework"]
    if "materials" in content:
        fields["materials"] = content["materials"]

    try:
        cfg = load_config(args.config)
    except (FileNotFoundError, RuntimeError, ValueError) as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1
    teacher_name, approved_names = resolve_allowlist(
        md, cfg, cli_allow_names=args.allow_names or "",
    )

    # Pre-flight scan: plan text + content JSON.
    pii_error = scan_for_pii(
        md + "\n" + json.dumps(content, default=str),
        teacher_name=teacher_name,
        approved_names=approved_names,
    )
    if pii_error:
        print(
            "Error: PII check failed. Remove student-identifying or sensitive information and retry.",
            file=sys.stderr,
        )
        return 1

    try:
        output_path = resolve_output_path(
            args.output, allow_anywhere=args.allow_anywhere
        )
    except ValueError as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1

    # Build the presentation in memory first so we can scan before writing.
    if args.template:
        tpl = Path(args.template)
        if tpl.exists():
            prs = _build_from_template(fields, tpl)
            mode_note = " (from template)"
        else:
            print(f"Warning: template not found ({tpl}), generating from scratch.",
                  file=sys.stderr)
            prs = _build_from_scratch(fields)
            mode_note = ""
    else:
        prs = _build_from_scratch(fields)
        mode_note = ""

    unresolved = find_unresolved_pptx_placeholders(prs)
    if unresolved:
        print(
            "Error: unresolved template placeholders remain after fill: "
            + ", ".join(unresolved[:10]),
            file=sys.stderr,
        )
        return 1

    # Post-fill scan: catches boilerplate in a district template and any
    # render-time surprises (e.g. names that appeared only after placeholder
    # substitution). Fail-closed to honor the FERPA claim.
    pii_error = scan_pptx_for_pii(
        prs, teacher_name=teacher_name, approved_names=approved_names,
    )
    if pii_error:
        print(
            "Error: PII check failed on rendered slide. Remove student-identifying or sensitive information and retry.",
            file=sys.stderr,
        )
        return 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"Wrote{mode_note}: {output_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
