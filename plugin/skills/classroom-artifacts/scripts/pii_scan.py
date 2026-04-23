#!/usr/bin/env python3
"""Lightweight PII scanner for classroom-artifacts."""
from __future__ import annotations

import locale
import re
import sys
from pathlib import Path
from typing import Any, Iterable, Optional

_SHARED_DIR = Path(__file__).resolve().parents[3] / "shared"
if str(_SHARED_DIR) not in sys.path:
    sys.path.insert(0, str(_SHARED_DIR))

try:
    from runtime_bootstrap import ensure_plugin_runtime_or_exit
except ImportError:  # pragma: no cover - exercised by isolated script tests
    ensure_plugin_runtime_or_exit = None

if ensure_plugin_runtime_or_exit is not None:
    ensure_plugin_runtime_or_exit(__file__)

from docx.text.paragraph import Paragraph

from pii_common import normalize_scan_text, scan_text_for_pii_matches

_IEP = re.compile(r"\b(IEP|BIP|504 plan|IEP goal|accommodations for \w+)\b", re.IGNORECASE)
PLACEHOLDER_PATTERN = re.compile(r"\{\{[A-Z0-9_]+\}\}")


def _read_text_file(path: Path, *, label: str) -> str:
    try:
        raw = path.read_bytes()
    except OSError as exc:
        raise ValueError(f"Could not read {label}: {exc}") from exc

    encodings: list[str] = ["utf-8", "utf-8-sig"]
    preferred = locale.getpreferredencoding(False)
    if preferred:
        normalized = preferred.lower().replace("_", "-")
        if normalized not in {"utf-8", "utf8", "utf-8-sig"}:
            encodings.append(preferred)

    last_error: UnicodeDecodeError | None = None
    for encoding in encodings:
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError as exc:
            last_error = exc

    encoding_list = ", ".join(encodings)
    raise ValueError(
        f"Could not decode {label}: {path}. Tried {encoding_list}."
    ) from last_error


# ---- Public API -----------------------------------------------------------

def extract_teacher_from_plan(md: str) -> str:
    """Find 'Teacher: <name>' in the plan markdown; returns '' if missing."""
    for line in md.splitlines():
        m = re.match(r"^\s*Teacher:\s*(.+?)\s*$", line)
        if m:
            return m.group(1).strip()
    return ""


def resolve_allowlist(plan_md, config, cli_allow_names=""):
    """Assemble the PII allowlist the way fill_template.py does.

    Returns (teacher_name, approved_names_list).
    """
    teacher = ""
    approved = []
    cfg = config or {}

    teacher_cfg = cfg.get("teacher") or {}
    if isinstance(teacher_cfg, dict):
        teacher = (teacher_cfg.get("name") or "").strip()

    # approved_names is a TOP-LEVEL list per shared/config-schema.md.
    for n in cfg.get("approved_names") or []:
        if isinstance(n, str):
            s = n.strip()
            if s:
                approved.append(s)

    for subj in cfg.get("subjects") or []:
        if isinstance(subj, dict):
            co = subj.get("co_teacher_name")
            if isinstance(co, str) and co.strip():
                approved.append(co.strip())

    plan_teacher = extract_teacher_from_plan(plan_md)
    if plan_teacher and not teacher:
        teacher = plan_teacher
    elif plan_teacher:
        approved.append(plan_teacher)

    if cli_allow_names:
        approved.extend(n.strip() for n in cli_allow_names.split(",") if n.strip())

    return teacher, approved


def load_plan_md(plan_path: Path) -> str:
    """Return markdown text for a plan argument.

    Accepts a .md / .markdown / .txt file directly, or a .docx file —
    in which case the sibling ``<stem>.plan.md`` sidecar (written by
    fill_template.py) is used. Raises FileNotFoundError with an
    actionable message if the plan cannot be resolved to markdown.
    """
    if not plan_path.exists():
        raise FileNotFoundError(f"plan not found: {plan_path}")

    suffix = plan_path.suffix.lower()
    if suffix in (".md", ".markdown", ".txt"):
        return _read_text_file(plan_path, label="plan markdown")

    if suffix == ".docx":
        sidecar = plan_path.with_suffix(".plan.md")
        if sidecar.exists():
            if sidecar.stat().st_mtime_ns < plan_path.stat().st_mtime_ns:
                raise FileNotFoundError(
                    f"Stale sidecar markdown found for {plan_path}: "
                    f"{sidecar} is older than the .docx. Re-run fill_template.py "
                    f"or pass the fresh .plan.md path to --plan."
                )
            return _read_text_file(sidecar, label="plan sidecar markdown")
        raise FileNotFoundError(
            f"No sidecar markdown found for {plan_path}. "
            f"Expected {sidecar}. Re-run fill_template.py (sidecars are "
            f"written by default) or pass the .plan.md path to --plan."
        )

    raise ValueError(
        f"Unsupported plan format '{suffix}'. Use .md or .docx (with sidecar)."
    )


def _iter_docx_text(doc) -> Iterable[str]:
    """Yield paragraph-level visible text from a python-docx Document.

    Covers the main body AND headers/footers on every section. Text
    boxes and nested tables are reached via ``w:p`` descent inside each
    root element — ``iter(w:p)`` walks into ``w:txbxContent`` wrappers
    too. Without this, template boilerplate in a header or sidebar
    textbox slipped past the "scan every draft" guarantee.
    """
    for para in _iter_docx_paragraphs(doc):
        text = "".join(run.text for run in para.runs if run.text)
        if text:
            yield text


def _iter_docx_paragraphs(doc) -> Iterable[Paragraph]:
    from docx.oxml.ns import qn

    roots: list[tuple[object, object]] = []
    body = getattr(getattr(doc, "element", None), "body", None)
    if body is not None:
        roots.append((body, doc))

    seen_ids: set[int] = set()
    for section in getattr(doc, "sections", []):
        for attr in (
            "header", "first_page_header", "even_page_header",
            "footer", "first_page_footer", "even_page_footer",
        ):
            hdr_ftr = getattr(section, attr, None)
            if hdr_ftr is None:
                continue
            element = getattr(hdr_ftr, "_element", None)
            if element is None:
                part = getattr(hdr_ftr, "part", None)
                element = getattr(part, "element", None) if part is not None else None
            if element is None or id(element) in seen_ids:
                continue
            seen_ids.add(id(element))
            roots.append((element, hdr_ftr))

    for root, parent in roots:
        for p in list(root.iter(qn("w:p"))):
            yield Paragraph(p, parent)


def _replace_in_text_runs(runs, replacements: dict[str, str]) -> bool:
    text_runs = [run for run in runs if run.text]
    if not text_runs:
        return False

    combined = "".join(run.text for run in text_runs)
    changed = False
    for key, value in replacements.items():
        if key in combined:
            combined = combined.replace(key, str(value))
            changed = True

    if changed:
        text_runs[0].text = combined
        for run in text_runs[1:]:
            run.text = ""
    return changed


def replace_docx_placeholders(doc, replacements: dict[str, str]) -> bool:
    changed = False
    for para in _iter_docx_paragraphs(doc):
        if _replace_in_text_runs(para.runs, replacements):
            changed = True
    return changed


def find_unresolved_docx_placeholders(doc) -> list[str]:
    found: set[str] = set()
    for para in _iter_docx_paragraphs(doc):
        text = "".join(run.text for run in para.runs if run.text)
        found.update(PLACEHOLDER_PATTERN.findall(text))
    return sorted(found)


def _iter_pptx_text(prs) -> Iterable[str]:
    """Yield paragraph-level visible text from a python-pptx Presentation.

    Recurses into grouped shapes, iterates table cells, and picks up
    speaker notes. Pre-fix, only top-level ``shape.text_frame`` text
    was scanned, so PII hidden inside a grouped shape, table cell, or
    the notes pane silently shipped to disk.
    """
    for para in _iter_pptx_paragraphs(prs):
        text = "".join(run.text for run in para.runs)
        if text:
            yield text


def _iter_pptx_paragraphs(prs):
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        group_type = MSO_SHAPE_TYPE.GROUP
    except Exception:
        group_type = None

    def walk(shapes):
        for shape in shapes:
            shape_type = getattr(shape, "shape_type", None)
            if group_type is not None and shape_type == group_type:
                yield from walk(shape.shapes)
                continue
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            yield para
                continue
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    yield para

    for slide in getattr(prs, "slides", []):
        yield from walk(slide.shapes)
        if getattr(slide, "has_notes_slide", False):
            notes_slide = slide.notes_slide
            notes_tf = getattr(notes_slide, "notes_text_frame", None)
            if notes_tf is not None:
                for para in notes_tf.paragraphs:
                    yield para


def replace_pptx_placeholders(prs, replacements: dict[str, str]) -> bool:
    changed = False
    for para in _iter_pptx_paragraphs(prs):
        if _replace_in_text_runs(para.runs, replacements):
            changed = True
    return changed


def find_unresolved_pptx_placeholders(prs) -> list[str]:
    found: set[str] = set()
    for para in _iter_pptx_paragraphs(prs):
        text = "".join(run.text for run in para.runs)
        found.update(PLACEHOLDER_PATTERN.findall(text))
    return sorted(found)


def scan_docx_for_pii(doc, teacher_name: str = "",
                      approved_names: Optional[list] = None) -> Optional[str]:
    """Scan a python-docx Document's visible text for PII. Returns error or None."""
    return scan_for_pii(
        "\n".join(_iter_docx_text(doc)),
        teacher_name=teacher_name,
        approved_names=approved_names,
    )


def scan_pptx_for_pii(prs, teacher_name: str = "",
                      approved_names: Optional[list] = None) -> Optional[str]:
    """Scan a python-pptx Presentation's visible text for PII. Returns error or None."""
    return scan_for_pii(
        "\n".join(_iter_pptx_text(prs)),
        teacher_name=teacher_name,
        approved_names=approved_names,
    )


def _scrub_markdown_headings(text: str) -> str:
    return "\n".join(
        line
        for line in text.splitlines()
        if not line.lstrip().startswith("#")
    )


def _issue_from_match(matched_text: str, label: str) -> str:
    if label == "SSN":
        return "SSN-like pattern detected."
    if label == "phone":
        return "Phone-like pattern detected."
    if label == "email address":
        return "Email address detected — remove it before writing."
    if label == "student ID":
        return "Student ID-like pattern detected."
    if label == "DOB / birthday field":
        return "DOB / birthday field detected."
    if label == "home address field":
        return "Home-address field detected."
    if label == "parent/guardian contact field":
        return "Parent/guardian contact field detected."
    if label == "lunch-status field":
        return "Lunch-status phrasing detected."
    if label == "medical-note field":
        return "Medical-note phrasing detected."
    if label == "discipline-note field":
        return "Discipline-note phrasing detected."
    if label == "roster-like name":
        return (
            f"roster-like keyword + name detected ('{matched_text}') — "
            "the 'Student:'/'Learner:' prefix denotes student context and "
            "cannot be allowlisted."
        )
    if label == "named accommodation":
        return (
            f"named accommodation detected ('{matched_text}') — "
            "IEP/504/ELL/SPED + a name is non-allowlistable."
        )
    if label == "bare name (unallowlisted)":
        return (
            f"possible student name '{matched_text}' (not in allowlist). "
            "Use --allow-names to accept."
        )
    if label == "bare name (all-caps, unallowlisted)":
        return f"possible student name '{matched_text}' (all-caps, not in allowlist)."
    if label == "initial-form name":
        return f"possible initial-form name '{matched_text}' (not in allowlist)."
    return f"PII-like content detected ('{matched_text}')."


def scan_for_pii(text, teacher_name="", approved_names=None):
    """Return an error-message string if suspicious PII is found; else None.

    Reuses the planner's shared scanner but returns one human-readable string
    so the artifact generators can print and exit cleanly.
    """
    if approved_names is None:
        approved_names = []

    all_allowed = set(approved_names)
    if teacher_name:
        all_allowed.add(teacher_name)

    text = _scrub_markdown_headings(normalize_scan_text(text))
    matches = scan_text_for_pii_matches(text, all_allowed)
    issues: list[str] = []

    def add_issue(message: str) -> None:
        if message not in issues:
            issues.append(message)

    if _IEP.search(text):
        add_issue("IEP / 504-plan phrasing detected — verify this isn't student-specific.")
    for matched_text, label in matches:
        add_issue(_issue_from_match(matched_text, label))

    return "; ".join(issues) if issues else None
