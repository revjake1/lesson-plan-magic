"""Regression tests for classroom-artifacts helpers.

Covers:
- FERPA-safe PII fail-closed for every helper (including agenda slide).
- --plan contract: markdown path, .docx+sidecar, .docx without sidecar.
- Exit ticket: all questions rendered, no silent drops.
- Agenda slide: homework override lands in the homework field (not materials).
"""
from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
from pathlib import Path

import pytest
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import artifact_common
import content_schema
import pii_scan
import generate_agenda_slide
import generate_do_now
import generate_exit_ticket
import generate_sub_plan


SCRIPTS_DIR = (
    Path(__file__).parent.parent / "skills" / "classroom-artifacts" / "scripts"
)
ASSETS_DIR = (
    Path(__file__).parent.parent / "skills" / "classroom-artifacts" / "assets"
)
SHARED_DIR = Path(__file__).parent.parent / "shared"


SAMPLE_PLAN_MD = """Subject: Chemistry
Teacher: Mr. Hallman

## 2026-04-22 — Wednesday

### Learning Intention
I can apply Newton's second law to everyday motion.

### Success Criteria
- I can identify forces acting on an object.
- I can write F = ma for a given scenario.

### Agenda
1. Warm-up: force diagram review
2. Mini-lecture
3. Guided practice

### Materials
- Worksheet

### Do Now
Sketch yesterday's lab.
"""


# ---------------------------------------------------------------------------
# Shared helpers: load_plan_md contract
# ---------------------------------------------------------------------------

class TestLoadPlanMd:
    def test_reads_md_directly(self, tmp_path):
        p = tmp_path / "plan.md"
        p.write_text(SAMPLE_PLAN_MD)
        assert pii_scan.load_plan_md(p) == SAMPLE_PLAN_MD

    def test_reads_markdown_extension(self, tmp_path):
        p = tmp_path / "plan.markdown"
        p.write_text(SAMPLE_PLAN_MD)
        assert pii_scan.load_plan_md(p) == SAMPLE_PLAN_MD

    def test_docx_with_sidecar(self, tmp_path):
        docx_path = tmp_path / "plan.docx"
        Document().save(str(docx_path))
        sidecar = tmp_path / "plan.plan.md"
        sidecar.write_text(SAMPLE_PLAN_MD)
        assert pii_scan.load_plan_md(docx_path) == SAMPLE_PLAN_MD

    def test_docx_with_stale_sidecar_errors(self, tmp_path):
        docx_path = tmp_path / "plan.docx"
        Document().save(str(docx_path))
        sidecar = tmp_path / "plan.plan.md"
        sidecar.write_text(SAMPLE_PLAN_MD)
        stale_ns = sidecar.stat().st_mtime_ns + 1_000_000
        os.utime(docx_path, ns=(stale_ns, stale_ns))

        with pytest.raises(FileNotFoundError, match="Stale sidecar"):
            pii_scan.load_plan_md(docx_path)

    def test_docx_without_sidecar_errors(self, tmp_path):
        docx_path = tmp_path / "plan.docx"
        Document().save(str(docx_path))
        with pytest.raises(FileNotFoundError):
            pii_scan.load_plan_md(docx_path)

    def test_missing_file_errors(self, tmp_path):
        with pytest.raises(FileNotFoundError):
            pii_scan.load_plan_md(tmp_path / "does_not_exist.md")

    def test_unsupported_extension_errors(self, tmp_path):
        p = tmp_path / "plan.pdf"
        p.write_text("not used")
        with pytest.raises(ValueError):
            pii_scan.load_plan_md(p)


class TestStructuredPlanLoading:
    def test_prefers_json_sidecar_for_docx_plan(self, tmp_path):
        docx_path = tmp_path / "plan.docx"
        Document().save(str(docx_path))
        markdown = SAMPLE_PLAN_MD.replace("2026-04-22", "2026-4-22")
        (tmp_path / "plan.plan.md").write_text(markdown)
        (tmp_path / "plan.plan.json").write_text(
            json.dumps(
                {
                    "v": 1,
                    "s": "JSON Chemistry",
                    "t": "Mr. Hallman",
                    "d": [
                        {
                            "dt": "2026-04-22",
                            "n": "Wednesday",
                            "li": "JSON learning intention.",
                            "sc": ["JSON criterion"],
                            "ag": ["JSON agenda"],
                            "m": ["JSON worksheet"],
                            "do": "JSON do now.",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )

        day = artifact_common.load_structured_day(
            docx_path, "2026-04-22", plan_md=markdown
        )

        assert day is not None
        assert day["subject"] == "JSON Chemistry"
        assert day["learning_intention"] == "JSON learning intention."
        assert day["agenda"] == ["JSON agenda"]

    def test_stale_json_sidecar_falls_back_to_markdown(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD, encoding="utf-8")
        json_sidecar = tmp_path / "plan.json"
        json_sidecar.write_text(
            json.dumps(
                {
                    "v": 1,
                    "s": "Stale JSON Subject",
                    "d": [{"dt": "2026-04-22", "li": "Stale JSON LI"}],
                }
            ),
            encoding="utf-8",
        )
        stale_ns = plan.stat().st_mtime_ns + 1_000_000
        os.utime(plan, ns=(stale_ns, stale_ns))

        day = artifact_common.load_structured_day(
            plan, "2026-04-22", plan_md=SAMPLE_PLAN_MD
        )

        assert day is not None
        assert day["subject"] == "Chemistry"
        assert day["learning_intention"] == (
            "I can apply Newton's second law to everyday motion."
        )


# ---------------------------------------------------------------------------
# Agenda slide — PII + homework override
# ---------------------------------------------------------------------------

def _run_helper(
    script_name,
    plan_path,
    out_path,
    *,
    env=None,
    allow_anywhere=True,
    **kwargs,
):
    argv = [
        sys.executable,
        str(SCRIPTS_DIR / script_name),
        "--plan", str(plan_path),
        "--date", "2026-04-22",
        "--subject", "chem",
        "--output", str(out_path),
    ]
    if allow_anywhere:
        argv.append("--allow-anywhere")
    for k, v in kwargs.items():
        flag = f"--{k.replace('_', '-')}"
        if isinstance(v, bool):
            if v:
                argv.append(flag)
            continue
        argv.extend([flag, str(v)])
    run_env = os.environ.copy()
    if env:
        run_env.update(env)
    return subprocess.run(argv, capture_output=True, text=True, env=run_env)


def _run_agenda(plan_path, out_path, *, env=None, allow_anywhere=True, **kwargs):
    return _run_helper(
        "generate_agenda_slide.py",
        plan_path,
        out_path,
        env=env,
        allow_anywhere=allow_anywhere,
        **kwargs,
    )


def _run_exit(plan_path, out_path, *, env=None, allow_anywhere=True, **kwargs):
    return _run_helper(
        "generate_exit_ticket.py",
        plan_path,
        out_path,
        env=env,
        allow_anywhere=allow_anywhere,
        **kwargs,
    )


def _run_do_now(plan_path, out_path, *, env=None, allow_anywhere=True, **kwargs):
    return _run_helper(
        "generate_do_now.py",
        plan_path,
        out_path,
        env=env,
        allow_anywhere=allow_anywhere,
        **kwargs,
    )


def _run_sub(plan_path, out_path, *, env=None, allow_anywhere=True, **kwargs):
    return _run_helper(
        "generate_sub_plan.py",
        plan_path,
        out_path,
        env=env,
        allow_anywhere=allow_anywhere,
        **kwargs,
    )


def _oversized_content_payload(runner) -> str:
    blob = "x" * ((1 * 1024 * 1024) + 128)
    if runner is _run_agenda:
        return json.dumps({"agenda": [blob]})
    if runner is _run_exit:
        return json.dumps({"questions": [blob]})
    if runner is _run_do_now:
        return json.dumps({"prompt": blob})
    if runner is _run_sub:
        return json.dumps({"backup": blob})
    raise AssertionError(f"Unexpected runner: {runner}")


def _pptx_text(path):
    prs = Presentation(str(path))
    out = []

    def walk(shapes):
        for shape in shapes:
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            out.append("".join(r.text for r in p.runs))
                continue
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)
                continue
            if getattr(shape, "has_text_frame", False):
                for p in shape.text_frame.paragraphs:
                    out.append("".join(r.text for r in p.runs))

    for slide in prs.slides:
        walk(slide.shapes)
    return out


def _docx_text(path):
    doc = Document(str(path))
    return [p.text for p in doc.paragraphs]


def _docx_all_text(path):
    doc = Document(str(path))
    out = [p.text for p in doc.paragraphs if p.text]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if para.text:
                        out.append(para.text)
    return out


class TestAgendaSlide:
    def test_ferpa_pii_blocks_output(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(
            "Subject: Chemistry\nTeacher: Mr. Hallman\n\n"
            "## 2026-04-22 — Wednesday\n\n"
            "### Learning Intention\nStudent: John Smith sets the pace today.\n"
            "\n### Success Criteria\n- Identify forces.\n"
            "\n### Agenda\n1. Warm-up\n"
        )
        out = tmp_path / "agenda.pptx"
        result = _run_agenda(plan, out)
        assert result.returncode != 0
        assert "PII check failed" in result.stderr
        assert not out.exists()

    def test_ferpa_ssn_blocks_output(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD + "\nParent SSN 123-45-6789\n")
        out = tmp_path / "agenda.pptx"
        result = _run_agenda(plan, out)
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()

    def test_homework_override_renders_in_homework_section(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "agenda.pptx"
        result = _run_agenda(
            plan, out,
            content_json=json.dumps({"homework": "Read chapter 3"}),
        )
        assert result.returncode == 0, result.stderr
        text = "\n".join(_pptx_text(out))
        assert "Homework" in text
        assert "Read chapter 3" in text
        # Regression: homework must NOT be dumped into the materials bullet list
        # (it used to overwrite fields['materials']).
        assert "• Read chapter 3" not in text

    def test_homework_override_in_template_fills_homework_placeholder(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "agenda.pptx"
        result = _run_agenda(
            plan, out,
            template=str(ASSETS_DIR / "agenda-slide-template.pptx"),
            content_json=json.dumps({"homework": "Read chapter 3"}),
        )
        assert result.returncode == 0, result.stderr
        text = "\n".join(_pptx_text(out))
        assert "Read chapter 3" in text
        assert "{{HOMEWORK}}" not in text

    def test_accepts_docx_plan_via_sidecar(self, tmp_path):
        docx_plan = tmp_path / "plan.docx"
        Document().save(str(docx_plan))
        (tmp_path / "plan.plan.md").write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "agenda.pptx"
        result = _run_agenda(docx_plan, out)
        assert result.returncode == 0, result.stderr
        assert out.exists()

    def test_docx_without_sidecar_errors_clearly(self, tmp_path):
        docx_plan = tmp_path / "plan.docx"
        Document().save(str(docx_plan))
        out = tmp_path / "agenda.pptx"
        result = _run_agenda(docx_plan, out)
        assert result.returncode != 0
        assert "sidecar" in result.stderr.lower()

    def test_non_padded_day_header_errors_loudly(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD.replace("2026-04-22", "2026-4-22"))
        out = tmp_path / "agenda.pptx"
        result = _run_agenda(plan, out)
        assert result.returncode != 0
        assert "zero-padded YYYY-MM-DD" in result.stderr
        assert not out.exists()

    def test_template_fills_table_placeholder(self, tmp_path):
        from pptx.util import Inches

        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)

        template = tmp_path / "agenda-template.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(0.7)).text_frame.text = (
            "{{LEARNING_INTENTION}}"
        )
        table = slide.shapes.add_table(1, 1, Inches(0.5), Inches(1.5), Inches(5), Inches(2)).table
        table.cell(0, 0).text = "{{AGENDA}}"
        prs.save(str(template))

        out = tmp_path / "agenda.pptx"
        result = _run_agenda(plan, out, template=str(template))

        assert result.returncode == 0, result.stderr
        text = "\n".join(_pptx_text(out))
        assert "I can apply Newton's second law to everyday motion." in text
        assert "1. Warm-up: force diagram review" in text
        assert "2. Mini-lecture" in text
        assert "{{AGENDA}}" not in text

    def test_unresolved_template_placeholder_blocks_output(self, tmp_path):
        from pptx.util import Inches

        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)

        template = tmp_path / "agenda-template.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(0.7)).text_frame.text = (
            "{{LEARNING_INTENTION}}"
        )
        slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(0.7)).text_frame.text = (
            "{{UNSUPPORTED_TOKEN}}"
        )
        prs.save(str(template))

        out = tmp_path / "agenda.pptx"
        result = _run_agenda(plan, out, template=str(template))

        assert result.returncode != 0
        assert "unresolved template placeholders" in result.stderr.lower()
        assert not out.exists()

    @pytest.mark.parametrize(
        ("runner", "out_name"),
        [
            (_run_agenda, "agenda.pptx"),
            (_run_exit, "exit.docx"),
            (_run_do_now, "do-now.docx"),
            (_run_sub, "sub-plan.docx"),
        ],
    )
    def test_helpers_prefer_json_sidecar_when_markdown_parse_would_fail(
        self, runner, out_name, tmp_path
    ):
        plan = tmp_path / "plan.md"
        plan.write_text(
            SAMPLE_PLAN_MD.replace("2026-04-22", "2026-4-22"),
            encoding="utf-8",
        )
        plan.with_suffix(".json").write_text(
            json.dumps(
                {
                    "v": 1,
                    "s": "Chemistry",
                    "t": "Mr. Hallman",
                    "d": [
                        {
                            "dt": "2026-04-22",
                            "n": "Wednesday",
                            "li": "JSON learning intention.",
                            "sc": ["JSON criterion"],
                            "ag": ["JSON agenda"],
                            "m": ["JSON worksheet"],
                            "do": "JSON do now.",
                        }
                    ],
                }
            ),
            encoding="utf-8",
        )
        out = tmp_path / out_name

        result = runner(plan, out)

        assert result.returncode == 0, result.stderr
        assert out.exists()


# ---------------------------------------------------------------------------
# Exit ticket — all questions, PII, .docx contract
# ---------------------------------------------------------------------------

class TestExitTicket:
    def test_all_questions_render_scratch_mode(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "exit.docx"
        result = _run_exit(
            plan, out,
            content_json=json.dumps(
                {"questions": ["Q1?", "Q2?", "Q3?", "Q4?", "Q5?"]}
            ),
        )
        assert result.returncode == 0, result.stderr
        text = "\n".join(_docx_text(out))
        for q in ("Q1?", "Q2?", "Q3?", "Q4?", "Q5?"):
            assert q in text, f"{q} missing from rendered ticket"

    def test_all_questions_render_template_mode_overflow(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "exit.docx"
        result = _run_exit(
            plan, out,
            template=str(ASSETS_DIR / "exit-ticket-template.docx"),
            content_json=json.dumps(
                {"questions": ["Q1?", "Q2?", "Q3?", "Q4?", "Q5?"]}
            ),
        )
        assert result.returncode == 0, result.stderr
        text = "\n".join(_docx_text(out))
        # Q1/Q2 fill the template placeholders; Q3-Q5 must appear in overflow.
        for q in ("Q1?", "Q2?", "Q3?", "Q4?", "Q5?"):
            assert q in text, f"{q} missing"

    def test_ferpa_pii_blocks_output(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "exit.docx"
        result = _run_exit(
            plan, out,
            content_json=json.dumps(
                {"questions": ["Explain why Taylor Johnson was correct."]}
            ),
        )
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()

    def test_docx_plan_without_sidecar_errors_not_unicode_crash(self, tmp_path):
        """Regression: pre-fix, `.docx` plans crashed with UnicodeDecodeError."""
        docx_plan = tmp_path / "plan.docx"
        Document().save(str(docx_plan))
        out = tmp_path / "exit.docx"
        result = _run_exit(docx_plan, out)
        assert result.returncode != 0
        # Must be the friendly sidecar message, NOT a UnicodeDecodeError traceback.
        assert "UnicodeDecodeError" not in result.stderr
        assert "sidecar" in result.stderr.lower()

    def test_non_padded_day_header_errors_loudly(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD.replace("2026-04-22", "2026-4-22"))
        out = tmp_path / "exit.docx"
        result = _run_exit(plan, out)
        assert result.returncode != 0
        assert "zero-padded YYYY-MM-DD" in result.stderr
        assert not out.exists()

    def test_docx_plan_with_sidecar_works(self, tmp_path):
        docx_plan = tmp_path / "plan.docx"
        Document().save(str(docx_plan))
        (tmp_path / "plan.plan.md").write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "exit.docx"
        result = _run_exit(
            docx_plan, out,
            content_json=json.dumps({"questions": ["Q1?"]}),
        )
        assert result.returncode == 0, result.stderr
        assert out.exists()

    def test_txt_output_writes_plain_text_not_ooxml(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "exit.txt"
        result = _run_exit(
            plan, out,
            content_json=json.dumps(
                {"questions": ["Q1?", "Q2?"], "metacognitive": "What still feels confusing?"}
            ),
        )

        assert result.returncode == 0, result.stderr
        raw = out.read_bytes()
        assert not raw.startswith(b"PK"), "plain-text mode wrote an OOXML zip"
        text = raw.decode("utf-8")
        assert "Exit Ticket" in text
        assert "1. Q1?" in text
        assert "2. Q2?" in text
        assert "What still feels confusing?" in text

    def test_txt_output_escapes_formula_prefixed_lines(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "exit.txt"
        result = _run_exit(
            plan,
            out,
            content_json=json.dumps(
                {
                    "learning_intention": '=HYPERLINK("https://example.com")',
                    "questions": ["Q1?"],
                    "metacognitive": "@everyone reflect on one confusion.",
                }
            ),
        )

        assert result.returncode == 0, result.stderr
        text = out.read_text(encoding="utf-8")
        assert '\'=HYPERLINK("https://example.com")' in text
        assert "'@everyone reflect on one confusion." in text

    def test_txt_output_rejects_template_argument(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "exit.txt"
        result = _run_exit(
            plan, out,
            template=str(ASSETS_DIR / "exit-ticket-template.docx"),
            content_json=json.dumps({"questions": ["Q1?"]}),
        )

        assert result.returncode != 0
        assert "templates are only supported for .docx" in result.stderr
        assert not out.exists()


# ---------------------------------------------------------------------------
# Do-now and sub plan — PII + .docx contract smoke tests
# ---------------------------------------------------------------------------

class TestDoNow:
    def test_ferpa_pii_blocks_output(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "do-now.docx"
        result = _run_do_now(
            plan, out,
            content_json=json.dumps(
                {"prompt": "Ask Taylor Johnson for her notes."}
            ),
        )
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()

    def test_docx_plan_without_sidecar_errors(self, tmp_path):
        docx_plan = tmp_path / "plan.docx"
        Document().save(str(docx_plan))
        out = tmp_path / "do-now.docx"
        result = _run_do_now(docx_plan, out)
        assert result.returncode != 0
        assert "UnicodeDecodeError" not in result.stderr
        assert "sidecar" in result.stderr.lower()

    def test_happy_path(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "do-now.docx"
        result = _run_do_now(
            plan, out,
            content_json=json.dumps({"prompt": "Sketch the graph."}),
        )
        assert result.returncode == 0, result.stderr
        text = "\n".join(_docx_text(out))
        assert "Sketch the graph." in text

    def test_template_header_placeholder_is_filled(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        template = tmp_path / "do-now-template.docx"
        doc = Document()
        doc.add_paragraph("{{DO_NOW_PROMPT}}")
        doc.sections[0].header.paragraphs[0].text = "Header {{DATE}}"
        doc.save(str(template))

        out = tmp_path / "do-now.docx"
        result = _run_do_now(
            plan,
            out,
            template=str(template),
            content_json=json.dumps({"prompt": "Sketch the graph."}),
        )

        assert result.returncode == 0, result.stderr
        written = Document(str(out))
        assert written.sections[0].header.paragraphs[0].text == "Header 2026-04-22"
        assert "Sketch the graph." in "\n".join(_docx_text(out))

    def test_non_padded_day_header_errors_loudly(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD.replace("2026-04-22", "2026-4-22"))
        out = tmp_path / "do-now.docx"
        result = _run_do_now(plan, out)
        assert result.returncode != 0
        assert "Traceback" not in result.stderr
        assert "zero-padded YYYY-MM-DD" in result.stderr
        assert not out.exists()


class TestSubPlan:
    def test_ferpa_pii_blocks_output(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "sub.docx"
        result = _run_sub(
            plan, out,
            content_json=json.dumps(
                {"backup": "Have Taylor Johnson lead the review."}
            ),
        )
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()

    def test_docx_plan_without_sidecar_errors(self, tmp_path):
        docx_plan = tmp_path / "plan.docx"
        Document().save(str(docx_plan))
        out = tmp_path / "sub.docx"
        result = _run_sub(docx_plan, out)
        assert result.returncode != 0
        assert "UnicodeDecodeError" not in result.stderr
        assert "sidecar" in result.stderr.lower()

    def test_happy_path(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "sub.docx"
        result = _run_sub(plan, out)
        assert result.returncode == 0, result.stderr
        assert out.exists()

    def test_starter_template_falls_back_to_plan_content(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "sub-template.docx"
        result = _run_sub(
            plan,
            out,
            template=str(ASSETS_DIR / "sub-plan-template.docx"),
        )
        assert result.returncode == 0, result.stderr
        text = "\n".join(_docx_all_text(out))
        assert "Learning Focus: I can apply Newton's second law to everyday motion." in text
        assert "Warm-up: force diagram review" in text
        assert "Worksheet" in text

    def test_starter_template_renders_return_notes_and_emergency(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "sub-template.docx"
        result = _run_sub(
            plan,
            out,
            template=str(ASSETS_DIR / "sub-plan-template.docx"),
            content_json=json.dumps(
                {
                    "activity": {"steps": ["Do packet"]},
                    "materials": ["Stopwatch"],
                    "backup": "Read chapter 3",
                    "emergency": {
                        "nearby_teacher": "Room 204",
                        "office_extension": "x201",
                    },
                    "return_notes": "Leave student work on my desk.",
                }
            ),
        )
        assert result.returncode == 0, result.stderr
        text = "\n".join(_docx_all_text(out))
        assert "Front office ext. x201" in text
        assert "Nearest colleague: Room 204" in text
        assert "Leave student work on my desk." in text

    def test_non_padded_day_header_errors_loudly(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD.replace("2026-04-22", "2026-4-22"))
        out = tmp_path / "sub.docx"
        result = _run_sub(plan, out)
        assert result.returncode != 0
        assert "Traceback" not in result.stderr
        assert "zero-padded YYYY-MM-DD" in result.stderr
        assert not out.exists()


class TestOutputPathFencing:
    @pytest.mark.parametrize(
        ("runner", "output_name", "kwargs"),
        [
            (_run_agenda, "agenda.pptx", {}),
            (_run_exit, "exit.docx", {"content_json": json.dumps({"questions": ["Q1?"]})}),
            (_run_do_now, "do-now.docx", {"content_json": json.dumps({"prompt": "Sketch the graph."})}),
            (_run_sub, "sub.docx", {}),
        ],
    )
    def test_blocks_output_outside_documented_dir_without_allow_anywhere(
        self, tmp_path, runner, output_name, kwargs
    ):
        home = tmp_path / "home"
        env = {"HOME": str(home)}
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / output_name

        result = runner(
            plan,
            out,
            env=env,
            allow_anywhere=False,
            **kwargs,
        )

        expected_root = home / "Documents" / "Lesson Plan Magic" / "outputs"
        assert result.returncode != 0
        assert str(expected_root) in result.stderr
        assert "allow-anywhere" in result.stderr.lower()
        assert not out.exists()

    @pytest.mark.parametrize(
        ("runner", "output_name", "kwargs"),
        [
            (_run_agenda, "agenda.pptx", {}),
            (_run_exit, "exit.docx", {"content_json": json.dumps({"questions": ["Q1?"]})}),
            (_run_do_now, "do-now.docx", {"content_json": json.dumps({"prompt": "Sketch the graph."})}),
            (_run_sub, "sub.docx", {}),
        ],
    )
    def test_allows_output_inside_documented_dir_by_default(
        self, tmp_path, runner, output_name, kwargs
    ):
        home = tmp_path / "home"
        env = {"HOME": str(home)}
        output_root = home / "Documents" / "Lesson Plan Magic" / "outputs"
        output_root.mkdir(parents=True)

        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = output_root / output_name

        result = runner(
            plan,
            out,
            env=env,
            allow_anywhere=False,
            **kwargs,
        )

        assert result.returncode == 0, result.stderr
        assert out.exists()


class TestContentValidation:
    @pytest.mark.parametrize(
        ("runner", "output_name", "payload"),
        [
            (_run_agenda, "agenda.pptx", json.dumps({"surprise": "text"})),
            (_run_exit, "exit.docx", json.dumps({"surprise": "text"})),
            (_run_do_now, "do-now.docx", json.dumps({"surprise": "text"})),
            (_run_sub, "sub.docx", json.dumps({"surprise": "text"})),
        ],
    )
    def test_unknown_content_keys_fail_closed(
        self, tmp_path, runner, output_name, payload
    ):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / output_name

        result = runner(plan, out, content_json=payload)

        assert result.returncode != 0
        assert "invalid content json" in result.stderr.lower()
        assert "unknown content key" in result.stderr.lower()
        assert not out.exists()

    @pytest.mark.parametrize(
        ("runner", "output_name", "payload", "fragment"),
        [
            (_run_agenda, "agenda.pptx", json.dumps({"agenda": [1]}), "content.agenda[0]"),
            (_run_exit, "exit.docx", json.dumps({"questions": [1]}), "content.questions[0]"),
            (_run_do_now, "do-now.docx", json.dumps({"prompt": ["hi"]}), "content.prompt"),
            (
                _run_sub,
                "sub.docx",
                json.dumps({"activity": {"estimated_min": "40"}}),
                "content.activity.estimated_min",
            ),
        ],
    )
    def test_bad_content_types_fail_closed(
        self, tmp_path, runner, output_name, payload, fragment
    ):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / output_name

        result = runner(plan, out, content_json=payload)

        assert result.returncode != 0
        assert "invalid content json" in result.stderr.lower()
        assert fragment in result.stderr
        assert not out.exists()

    @pytest.mark.parametrize(
        ("runner", "output_name"),
        [
            (_run_agenda, "agenda.pptx"),
            (_run_exit, "exit.docx"),
            (_run_do_now, "do-now.docx"),
            (_run_sub, "sub.docx"),
        ],
    )
    def test_null_content_json_fails_closed(self, tmp_path, runner, output_name):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / output_name

        result = runner(plan, out, content_json="null")

        assert result.returncode != 0
        assert "invalid content json" in result.stderr.lower()
        assert "json null" in result.stderr.lower()
        assert not out.exists()

    @pytest.mark.parametrize(
        ("runner", "output_name"),
        [
            (_run_agenda, "agenda.pptx"),
            (_run_exit, "exit.docx"),
            (_run_do_now, "do-now.docx"),
            (_run_sub, "sub.docx"),
        ],
    )
    def test_oversized_content_json_fails_closed(self, tmp_path, runner, output_name):
        validator = {
            _run_agenda: content_schema.validate_agenda_content,
            _run_exit: content_schema.validate_exit_ticket_content,
            _run_do_now: content_schema.validate_do_now_content,
            _run_sub: content_schema.validate_sub_plan_content,
        }[runner]

        with pytest.raises(content_schema.ContentValidationError, match="content limit"):
            content_schema.load_and_validate_content(
                content_file=None,
                content_json=_oversized_content_payload(runner),
                validator=validator,
            )


class TestConfigLoading:
    @pytest.mark.parametrize(
        ("runner", "output_name", "kwargs"),
        [
            (_run_agenda, "agenda.pptx", {}),
            (_run_exit, "exit.docx", {"content_json": json.dumps({"questions": ["Q1?"]})}),
            (_run_do_now, "do-now.docx", {"content_json": json.dumps({"prompt": "Sketch the graph."})}),
            (_run_sub, "sub.docx", {}),
        ],
    )
    def test_malformed_config_fails_closed(self, tmp_path, runner, output_name, kwargs):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        config = tmp_path / "broken.yaml"
        config.write_text("approved_names: [123\n", encoding="utf-8")
        out = tmp_path / output_name

        result = runner(plan, out, config=str(config), **kwargs)

        assert result.returncode != 0
        assert "could not load config" in result.stderr.lower()
        assert not out.exists()

    @pytest.mark.parametrize(
        ("runner", "output_name", "kwargs"),
        [
            (_run_agenda, "agenda.pptx", {}),
            (_run_exit, "exit.docx", {"content_json": json.dumps({"questions": ["Q1?"]})}),
            (_run_do_now, "do-now.docx", {"content_json": json.dumps({"prompt": "Sketch the graph."})}),
            (_run_sub, "sub.docx", {}),
        ],
    )
    def test_non_mapping_config_fails_closed(self, tmp_path, runner, output_name, kwargs):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        config = tmp_path / "broken.yaml"
        config.write_text("- just\n- a\n- list\n", encoding="utf-8")
        out = tmp_path / output_name

        result = runner(plan, out, config=str(config), **kwargs)

        assert result.returncode != 0
        assert "yaml mapping" in result.stderr.lower()
        assert not out.exists()


# ---------------------------------------------------------------------------
# PII scanner helpers
# ---------------------------------------------------------------------------

class TestPiiScanner:
    """Regression tests for the scanner itself."""

    def test_sentence_starter_plus_name_is_flagged(self):
        """Pre-fix, 'Ask Taylor Johnson' slipped past because the
        sentence-starter branch always continued."""
        result = pii_scan.scan_for_pii("Ask Taylor Johnson for her notes.")
        assert result is not None and "Taylor" in result

    def test_triple_capitalized_sequence_catches_inner_pair(self):
        """Pre-fix, non-overlapping finditer consumed 'Today Student',
        leaving 'Student John' undetected."""
        result = pii_scan.scan_for_pii("Today Student John will lead the lab.")
        assert result is not None and "Student John" in result

    def test_resolve_allowlist_ignores_non_string_approved_names(self):
        teacher, approved = pii_scan.resolve_allowlist(
            SAMPLE_PLAN_MD,
            {"approved_names": [123, ["Jane Doe"], "Rosa Parks"]},
        )
        assert teacher == "Mr. Hallman"
        assert approved == ["Rosa Parks"]

    def test_structural_heading_if_students_passes(self):
        """'If Students Finish Early (Backup Plan)' is a template heading,
        not a name."""
        assert pii_scan.scan_for_pii(
            "If Students Finish Early (Backup Plan)"
        ) is None

    def test_default_allowlist_hero_names_pass(self):
        assert pii_scan.scan_for_pii("Albert Einstein derived the equation.") is None

    def test_non_ascii_name_is_flagged(self):
        result = pii_scan.scan_for_pii("José García needs support.")
        assert result is not None and "José García" in result

    def test_place_name_like_phrase_passes(self):
        assert pii_scan.scan_for_pii("Meet in Austin Clinic after school.") is None

    def test_ap_biology_passes(self):
        assert pii_scan.scan_for_pii("AP Biology lab practical on Friday.") is None

    def test_direct_ssn_detected(self):
        assert pii_scan.scan_for_pii("Parent SSN 123-45-6789") is not None

    def test_dashed_phone_detected(self):
        assert pii_scan.scan_for_pii("Call 555-555-5555 if absent.") is not None

    def test_dotted_phone_detected(self):
        """Pre-fix, _PHONE only matched dashes, so 555.555.5555 slipped past
        the artifact scan even though fill_template.py would block it."""
        assert pii_scan.scan_for_pii("Call 555.555.5555 if absent.") is not None

    def test_bare_ten_digit_phone_detected(self):
        """Pre-fix, 5555555555 was not caught by the artifact scanner."""
        assert pii_scan.scan_for_pii("Call 5555555555 if absent.") is not None

    def test_school_org_email_detected(self):
        """Pre-fix, firstname.lastname@school.org slipped past the artifact
        scanner even though fill_template.py rejects it — the artifact
        regex only caught ``student*@...`` or ``name.YYYY@...``."""
        assert pii_scan.scan_for_pii(
            "Email john.smith@school.org if you are absent."
        ) is not None

    def test_district_edu_email_detected(self):
        assert pii_scan.scan_for_pii("Contact jane@district.edu.") is not None

    def test_k12_email_detected(self):
        assert pii_scan.scan_for_pii(
            "Reach out to alex@elem.k12.ga.us with questions."
        ) is not None

    def test_student_id_detected_after_nfkc_normalization(self):
        result = pii_scan.scan_for_pii("Student ID: １２３４５６")
        assert result is not None and "Student ID" in result

    def test_dob_detected(self):
        result = pii_scan.scan_for_pii("DOB: 04/22/2012")
        assert result is not None and "DOB" in result

    def test_home_address_detected(self):
        result = pii_scan.scan_for_pii("Home address: 123 Main Street, Macon, GA")
        assert result is not None and "address" in result.lower()

    def test_parent_contact_detected(self):
        result = pii_scan.scan_for_pii("Parent contact: Maria Rodriguez")
        assert result is not None and "parent/guardian" in result.lower()

    def test_lunch_status_detected(self):
        result = pii_scan.scan_for_pii("Student receives free-and-reduced lunch.")
        assert result is not None and "lunch" in result.lower()

    def test_medical_note_detected(self):
        result = pii_scan.scan_for_pii("Medication note: uses inhaler at lunch.")
        assert result is not None and "medical" in result.lower()

    def test_discipline_note_detected(self):
        result = pii_scan.scan_for_pii("Discipline record: office referral from Monday.")
        assert result is not None and "discipline" in result.lower()

    def test_zero_width_joiners_are_stripped_before_scan(self):
        result = pii_scan.scan_for_pii("Call 555\u200b555\u200b5555 if absent.")
        assert result is not None and "Phone" in result

    def test_multiple_bare_names_are_all_reported(self):
        result = pii_scan.scan_for_pii("Johnny Doe partnered with Mary Jones.")
        assert result is not None
        assert "Johnny Doe" in result
        assert "Mary Jones" in result


class TestDoNowPhonePii:
    """End-to-end: phone numbers in content must block do-now output."""

    def test_dotted_phone_in_content_blocks_output(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "do-now.docx"
        result = _run_do_now(
            plan, out,
            content_json=json.dumps({"prompt": "Call 555.555.5555 if absent."}),
        )
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()

    def test_bare_ten_digit_phone_in_content_blocks_output(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "do-now.docx"
        result = _run_do_now(
            plan, out,
            content_json=json.dumps({"prompt": "Call 5555555555 if absent."}),
        )
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()

    def test_school_org_email_in_content_blocks_output(self, tmp_path):
        """End-to-end: a school-domain email in the do-now prompt must block
        output (was writing ``Email john.smith@school.org`` pre-fix)."""
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "do-now.docx"
        result = _run_do_now(
            plan, out,
            content_json=json.dumps(
                {"prompt": "Email john.smith@school.org if you are absent."}
            ),
        )
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()


class TestBrokenInstallFailsClosed:
    """Missing pii_scan.py must abort loudly — no fail-open stubs.

    Pre-fix, each helper wrapped ``from pii_scan import ...`` in
    ``except ImportError`` and substituted no-op stubs, so a broken
    install would still write student names to disk.
    """

    HELPERS = (
        "generate_do_now.py",
        "generate_exit_ticket.py",
        "generate_sub_plan.py",
        "generate_agenda_slide.py",
    )
    LOCAL_DEPS = (
        "artifact_common.py",
        "content_schema.py",
        "plan_parser.py",
    )

    @pytest.mark.parametrize("helper", HELPERS)
    def test_helper_without_pii_scan_errors(self, helper, tmp_path):
        for dep in self.LOCAL_DEPS:
            shutil.copy(SCRIPTS_DIR / dep, tmp_path / dep)
        dest = tmp_path / helper
        shutil.copy(SCRIPTS_DIR / helper, dest)
        # Intentionally DO NOT copy pii_scan.py.

        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out_name = "agenda.pptx" if helper.endswith("agenda_slide.py") else "out.docx"
        out = tmp_path / out_name

        # Scrub PYTHONPATH so the helper can't reach the real pii_scan via
        # inherited test-env paths — this simulates a broken/incomplete install.
        env = {k: v for k, v in os.environ.items() if k != "PYTHONPATH"}
        result = subprocess.run(
            [
                sys.executable, str(dest),
                "--plan", str(plan),
                "--date", "2026-04-22",
                "--subject", "chem",
                "--output", str(out),
            ],
            capture_output=True, text=True, env=env,
        )
        assert result.returncode != 0, (
            f"{helper} must abort loudly when pii_scan is missing; "
            f"stdout={result.stdout!r} stderr={result.stderr!r}"
        )
        assert not out.exists(), (
            f"{helper} wrote an artifact with PII checks disabled"
        )
        combined = result.stdout + result.stderr
        # Must surface the ImportError, not silently succeed.
        assert "pii_scan" in combined or "ImportError" in combined or "ModuleNotFoundError" in combined


class TestSelfContainedInstall:
    HELPERS = (
        (
            "generate_do_now.py",
            "out.docx",
            {"content_json": json.dumps({"prompt": "Sketch the graph."})},
        ),
        (
            "generate_exit_ticket.py",
            "out.docx",
            {"content_json": json.dumps({"questions": ["Q1?"]})},
        ),
        (
            "generate_sub_plan.py",
            "out.docx",
            {},
        ),
        (
            "generate_agenda_slide.py",
            "out.pptx",
            {},
        ),
    )
    LOCAL_DEPS = (
        "artifact_common.py",
        "content_schema.py",
        "pii_scan.py",
        "plan_parser.py",
    )

    @pytest.mark.parametrize(("helper", "out_name", "kwargs"), HELPERS)
    def test_helper_runs_without_lesson_planner_dependency(
        self, helper, out_name, kwargs, tmp_path
    ):
        for dep in self.LOCAL_DEPS:
            shutil.copy(SCRIPTS_DIR / dep, tmp_path / dep)
        shutil.copy(SHARED_DIR / "pii_common.py", tmp_path / "pii_common.py")
        dest = tmp_path / helper
        shutil.copy(SCRIPTS_DIR / helper, dest)

        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / out_name

        argv = [
            sys.executable,
            str(dest),
            "--plan", str(plan),
            "--date", "2026-04-22",
            "--subject", "chem",
            "--output", str(out),
            "--allow-anywhere",
        ]
        for key, value in kwargs.items():
            argv.extend([f"--{key.replace('_', '-')}", str(value)])

        env = {k: v for k, v in os.environ.items() if k != "PYTHONPATH"}
        result = subprocess.run(
            argv,
            capture_output=True,
            text=True,
            env=env,
            cwd=tmp_path,
        )

        assert result.returncode == 0, (
            f"{helper} should run from the classroom-artifacts bundle alone; "
            f"stdout={result.stdout!r} stderr={result.stderr!r}"
        )
        assert out.exists()


class TestPiiHelpers:
    def test_scan_docx_for_pii_catches_name_in_paragraph(self, tmp_path):
        doc = Document()
        doc.add_paragraph("Student response: Taylor Johnson")
        assert pii_scan.scan_docx_for_pii(doc) is not None

    def test_scan_docx_for_pii_clean_doc(self, tmp_path):
        doc = Document()
        doc.add_paragraph("Exit Ticket for Chemistry")
        doc.add_paragraph("Question 1")
        assert pii_scan.scan_docx_for_pii(doc) is None

    def test_scan_pptx_for_pii_catches_name(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx = slide.shapes.add_textbox(0, 0, 1000000, 500000)
        tx.text_frame.text = "Today's exemplar: Taylor Johnson"
        assert pii_scan.scan_pptx_for_pii(prs) is not None

    def test_scan_pptx_for_pii_clean_deck(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx = slide.shapes.add_textbox(0, 0, 1000000, 500000)
        tx.text_frame.text = "Today's Agenda"
        assert pii_scan.scan_pptx_for_pii(prs) is None

    def test_scan_pptx_catches_pii_in_table_cell(self, tmp_path):
        """Regression (P1): PII inside a PPTX table cell used to slip past
        — _iter_pptx_text only walked top-level ``shape.text_frame`` text.
        """
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table_shape = slide.shapes.add_table(2, 2, 0, 0, 4000000, 2000000)
        table_shape.table.cell(0, 0).text = "Name"
        table_shape.table.cell(1, 0).text = "Taylor Johnson"
        assert pii_scan.scan_pptx_for_pii(prs) is not None

    def test_scan_pptx_catches_pii_in_speaker_notes(self, tmp_path):
        """Regression (P1): speaker-notes text used to ship without scan."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        notes = slide.notes_slide.notes_text_frame
        notes.text = "Remind Taylor Johnson about the test."
        assert pii_scan.scan_pptx_for_pii(prs) is not None

    def test_scan_pptx_catches_pii_in_grouped_shape(self, tmp_path):
        """Regression (P1): shapes inside a PPTX group didn't get walked.

        We can't directly construct a group via python-pptx's high-level
        API, so we verify the walker handles the GROUP type by calling
        _iter_pptx_text on a presentation whose group is built via XML —
        or, simpler here, by patching a shape to report as a group.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        inner = slide.shapes.add_textbox(0, 0, 1000000, 500000)
        inner.text_frame.text = "Student: John Doe"

        class _FakeGroup:
            shape_type = MSO_SHAPE_TYPE.GROUP
            has_table = False
            has_text_frame = False
            shapes = [inner]

        class _FakeShapes:
            def __iter__(self_inner):
                return iter([_FakeGroup()])

        class _FakeSlide:
            shapes = _FakeShapes()
            has_notes_slide = False

        class _FakePrs:
            slides = [_FakeSlide()]

        assert pii_scan.scan_pptx_for_pii(_FakePrs()) is not None


# ---------------------------------------------------------------------------
# End-to-end bypass attempts (P0 / P1 regressions from the pre-release audit)
# ---------------------------------------------------------------------------

class TestAllowlistBypassClosed:
    """The P0 bug: ``--allow-names 'Jane Doe'`` + ``Student: Jane Doe`` in
    the plan or content wrote a student roster to disk. The fix: roster
    keyword context is non-allowlistable."""

    def test_agenda_blocks_allowlist_bypass_in_plan(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(
            "Subject: Chemistry\nTeacher: Mr. Hallman\n\n"
            "## 2026-04-22 — Wednesday\n\n"
            "### Learning Intention\nStudent: Jane Doe leads today.\n"
            "\n### Agenda\n1. Warm-up\n"
        )
        out = tmp_path / "agenda.pptx"
        result = _run_agenda(plan, out, allow_names="Jane Doe")
        assert result.returncode != 0, (
            "Agenda slide allowed --allow-names to bypass roster-keyword PII"
        )
        assert "PII" in result.stderr
        assert not out.exists()

    def test_do_now_blocks_allowlist_bypass_in_content(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "do-now.docx"
        result = _run_do_now(
            plan, out,
            allow_names="Jane Doe",
            content_json=json.dumps({"prompt": "Student: Jane Doe, report to office."}),
        )
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()

    def test_all_caps_roster_blocked(self, tmp_path):
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "do-now.docx"
        result = _run_do_now(
            plan, out,
            content_json=json.dumps({"prompt": "STUDENT: JOHN DOE, see me after class."}),
        )
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()

    def test_personal_email_domain_blocked(self, tmp_path):
        """Regression: personal-domain emails (gmail.com / yahoo.com)
        previously passed the artifact scan because the regex only
        matched .k12.*/.edu/.org."""
        plan = tmp_path / "plan.md"
        plan.write_text(SAMPLE_PLAN_MD)
        out = tmp_path / "do-now.docx"
        result = _run_do_now(
            plan, out,
            content_json=json.dumps({"prompt": "Email john.doe@gmail.com for the worksheet."}),
        )
        assert result.returncode != 0
        assert "PII" in result.stderr
        assert not out.exists()
